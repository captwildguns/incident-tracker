"""Build the PowerPoint export from the per-slide PNGs in raw/slides/.

Run shoot-slides.cjs first (it needs the deck served on port 8899), then this.
One image per slide, so the layout matches the HTML exactly. Text is not
editable in PowerPoint; the HTML deck is the editable source.

Presenter notes live in NOTES below, keyed by slide number. Keep them in sync
when slides are added or reordered.
"""
import glob
import os
import sys
from pptx import Presentation
from pptx.util import Inches

HERE = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
SLIDES = os.path.join(HERE, 'raw', 'slides', '*.png')
OUT = os.path.join(os.path.dirname(HERE), 'Incident-Tracker-Launch-Session-1.pptx')

NOTES = {
    1: "30 minutes, every Friday. Purpose: work the Incidents Launch Plan to the "
       "November 25, 2026 release.",
    2: "Today's plan, six subjects. Keep the room moving.",
    3: "BK opens here. The plan belongs to the people in the room: every item has a "
       "person on it, that person owns status, and anyone can add an item their area "
       "needs or call out an item that does not help the launch.",
    4: "Still BK: the shape of the plan. The workstreams are containers; the items inside "
       "them are what people own.",
    5: "Background. The gap Incidents fills, and that the design is settled.",
    6: "The process is part of what makes this project unusual: AI drove the design, the "
       "mockups, and the documents. It began as a Figma Make design and moved to Claude "
       "mid-project, where Claude Code built the working Forge prototype. Say AI-powered, "
       "not hackathon.",
    7: "SAY THIS: an incident almost never gets written up while it is happening. The "
       "driver's job in the moment is to keep 50 kids safe and keep the bus moving. The "
       "report comes later, once they have pulled into the school loop or gotten back to "
       "the garage. Note we are not showing the driver tablet today.",
    8: "The only system slide. Name the five capabilities, do not explain them here: the "
       "demo does that. Data is realistic but simulated, so no real notifications go out.",
    9: "Switch to captwildguns.github.io/incident-tracker. Five beats, about 15 minutes. "
       "This is where the detail lives now.",
    10: "Show incidents-vnext.staging.student-transportation.tylerapp.com briefly. Student "
        "names in this screenshot are blurred.",
    11: "Going forward. From next Friday the board gets filtered by person and each owner "
        "gives status on their own items.",
}


def main():
    src = sorted(glob.glob(SLIDES))
    if not src:
        sys.exit('No slide PNGs found. Run tools/shoot-slides.cjs first.')

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    for i, png in enumerate(src, start=1):
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(png, 0, 0, width=prs.slide_width, height=prs.slide_height)
        slide.notes_slide.notes_text_frame.text = NOTES.get(i, '')

    prs.save(OUT)
    missing = [n for n in range(1, len(src) + 1) if n not in NOTES]
    print('wrote %s: %d slides, %.2f MB' % (
        os.path.basename(OUT), len(src), os.path.getsize(OUT) / 1024 / 1024))
    if missing:
        print('WARNING: no presenter notes for slides', missing)


if __name__ == '__main__':
    main()
